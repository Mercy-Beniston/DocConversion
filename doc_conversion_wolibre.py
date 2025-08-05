import os
import boto3
import time
from PIL import Image
from fpdf import FPDF
from docx import Document
from pptx import Presentation

textract_client = boto3.client('textract', region_name='us-east-1')
s3_bucket = "your-s3-bucket"
s3_key = "uploads/converted.pdf"

input_file = "input_file_path"  # Update this
converted_pdf = "converted.pdf"
output_pdf = "ocr_output.pdf"

# Step 1: Extract Images from DOCX
def extract_images_from_docx(docx_path):
    doc = Document(docx_path)
    images = []
    for rel in doc.part._rels:
        rel_obj = doc.part._rels[rel]
        if "image" in rel_obj.target_ref:
            img_data = rel_obj.target_part.blob
            img_name = rel_obj.target_ref.split('/')[-1]
            with open(img_name, 'wb') as f:
                f.write(img_data)
            images.append(img_name)
    return images

# Step 2: Extract Images from PPTX
def extract_images_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    images = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                image = shape.image
                img_name = f"{image.sha1}.jpg"
                with open(img_name, 'wb') as f:
                    f.write(image.blob)
                images.append(img_name)
    return images

# Step 3: Convert Images to PDF
def convert_images_to_pdf(image_paths, output_pdf):
    images = [Image.open(img).convert('RGB') for img in image_paths]
    images[0].save(output_pdf, save_all=True, append_images=images[1:])
    for img in image_paths:
        os.remove(img)

# Step 4: Upload to S3
def upload_to_s3(file_path, bucket, key):
    s3 = boto3.client('s3')
    s3.upload_file(file_path, bucket, key)

# Step 5: Start Textract Job
def start_textract_job(bucket, key):
    response = textract_client.start_document_text_detection(
        DocumentLocation={'S3Object': {'Bucket': bucket, 'Name': key}}
    )
    return response['JobId']

# Step 6: Wait for Textract
def wait_for_job(job_id):
    while True:
        response = textract_client.get_document_text_detection(JobId=job_id)
        status = response['JobStatus']
        if status in ['SUCCEEDED', 'FAILED']:
            return response
        time.sleep(5)

# Step 7: Extract Text
def extract_text_from_textract(response):
    text = ""
    next_token = None
    while True:
        if next_token:
            response = textract_client.get_document_text_detection(JobId=response['JobId'], NextToken=next_token)
        blocks = response['Blocks']
        for block in blocks:
            if block['BlockType'] == 'LINE':
                text += block['Text'] + '\n'
        next_token = response.get('NextToken')
        if not next_token:
            break
    return text

# Step 8: Simple PDF with Text
def generate_pdf_with_text(text, output_pdf):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    pdf.multi_cell(0, 10, text)
    pdf.output(output_pdf)

# ======== Full Flow ========
ext = os.path.splitext(input_file)[1].lower()

if ext == ".docx":
    imgs = extract_images_from_docx(input_file)
    if not imgs:
        raise Exception("No images found, OCR only applies to scanned images.")
    convert_images_to_pdf(imgs, converted_pdf)

elif ext == ".pptx":
    imgs = extract_images_from_pptx(input_file)
    if not imgs:
        raise Exception("No images found, OCR only applies to scanned images.")
    convert_images_to_pdf(imgs, converted_pdf)

elif ext in [".jpg", ".jpeg", ".png", ".tiff", ".bmp"]:
    convert_images_to_pdf([input_file], converted_pdf)

elif ext == ".pdf":
    converted_pdf = input_file

else:
    raise Exception(f"Unsupported file type: {ext}")

# Textract OCR Process
upload_to_s3(converted_pdf, s3_bucket, s3_key)
job_id = start_textract_job(s3_bucket, s3_key)
print("Textract job started, waiting...")
response = wait_for_job(job_id)
print("Textract completed.")
extracted_text = extract_text_from_textract(response)
generate_pdf_with_text(extracted_text, output_pdf)
print(f"OCR Searchable PDF saved as {output_pdf}")
