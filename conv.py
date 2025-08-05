import os
import boto3
import time
from PIL import Image
from fpdf import FPDF
from docx import Document
from pptx import Presentation
import zipfile

boto3.setup_default_session(profile_name=os.getenv("profile_name"))
textract_client = boto3.client('textract', region_name='us-east-2')
#bedrock = boto3.client('bedrock-runtime', 'us-east-1', endpoint_url='https://bedrock-runtime.us-east-1.amazonaws.com',
                       #config=config)
s3_bucket = "doc-conv-poc"#https://doc-conv-poc.s3.us-east-2.amazonaws.com/input_do.pdf
s3_key = "uploads/converted_xlsx.pdf"
#input_file = "input_do.docx"
input_file = "sample_excel.xlsx"
converted_pdf = "converted_xlsx.pdf"
output_pdf = "searchable_text_only_xlsx.pdf"
s3_searchable_key = "uploads/searchable_text_only_xlsx.pdf"


# Extract images from DOCX
def extract_images_from_docx(docx_path):
    doc = Document(docx_path)
    images = []
    for rel in doc.part._rels:
        rel_obj = doc.part._rels[rel]
        if "image" in rel_obj.target_ref:
            img_data = rel_obj.target_part.blob
            img_name = rel_obj.target_ref.split('/')[-1]
            with open(img_name, 'wb') as f:
                f.write(img_data)
            images.append(img_name)
    return images


# Extract images from PPTX
def extract_images_from_pptx(pptx_path):
    print(pptx_path)
    prs = Presentation(pptx_path)
    images = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:
                image = shape.image
                img_name = f"{image.sha1}.jpg"
                with open(img_name, 'wb') as f:
                    f.write(image.blob)
                images.append(img_name)
    return images


# Extract images from XLSX (basic version)
def extract_images_from_xlsx(xlsx_path):
    images = []
    with zipfile.ZipFile(xlsx_path, 'r') as z:
        for file in z.namelist():
            if file.startswith('xl/media/') and (
                    file.endswith('.png') or file.endswith('.jpg') or file.endswith('.jpeg') or file.endswith(
                    '.bmp') or file.endswith('.gif')):
                img_data = z.read(file)
                img_name = os.path.basename(file)
                with open(img_name, 'wb') as f:
                    f.write(img_data)
                images.append(img_name)
    return images


# Convert images to PDF
def convert_images_to_pdf(image_paths, output_pdf):
    images = [Image.open(img).convert('RGB') for img in image_paths]
    images[0].save(output_pdf, save_all=True, append_images=images[1:])
    for img in image_paths:
        os.remove(img)


# Upload to S3
def upload_to_s3(file_path, bucket, key):
    s3 = boto3.client('s3')
    s3.upload_file(file_path, bucket, key)


# Start Textract async job
def start_textract_job(bucket, key):
    response = textract_client.start_document_text_detection(
        DocumentLocation={'S3Object': {'Bucket': bucket, 'Name': key}}
    )
    return response['JobId']


# Wait for Textract job completion
def wait_for_job(job_id):
    while True:
        response = textract_client.get_document_text_detection(JobId=job_id)
        status = response['JobStatus']
        if status in ['SUCCEEDED', 'FAILED']:
            return response
        time.sleep(5)


# Extract plain text from Textract results
def extract_text_from_textract(job_id):
    full_text = ""
    next_token = None
    while True:
        if next_token:
            response = textract_client.get_document_text_detection(JobId=job_id, NextToken=next_token)
        else:
            response = textract_client.get_document_text_detection(JobId=job_id)

        for block in response['Blocks']:
            if block['BlockType'] == 'LINE':
                full_text += block['Text'] + "\n"

        next_token = response.get('NextToken')
        if not next_token:
            break
    return full_text


# Generate plain text-only PDF
def generate_text_pdf(text, output_pdf):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    pdf.multi_cell(0, 10, text)
    pdf.output(output_pdf)
    print(f"Text-only searchable PDF saved as {output_pdf}")
    upload_to_s3(output_pdf, s3_bucket, s3_searchable_key)


# ===== Full Flow =====

ext = os.path.splitext(input_file)[1].lower()

if ext in [".docx"]:
    images = extract_images_from_docx(input_file)
elif ext in [".pptx"]:
    images = extract_images_from_pptx(input_file)
elif ext in [".xlsx"]:
    images = extract_images_from_xlsx(input_file)
elif ext in [".jpg", ".jpeg", ".png", ".tiff", ".bmp"]:
    images = [input_file]
else:
    raise Exception(f"Unsupported file type: {ext}")

if not images:
    raise Exception("No images found to process.")

convert_images_to_pdf(images, converted_pdf)
upload_to_s3(converted_pdf, s3_bucket, s3_key)

job_id = start_textract_job(s3_bucket, s3_key)
print(f"Textract job started: {job_id}")
response = wait_for_job(job_id)

if response['JobStatus'] != 'SUCCEEDED':
    raise Exception("Textract job failed.")

print("Textract completed. Generating searchable PDF...")
extracted_text = extract_text_from_textract(job_id)
generate_text_pdf(extracted_text, output_pdf)
